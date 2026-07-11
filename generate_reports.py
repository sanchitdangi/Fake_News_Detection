import os
import pandas as pd
import json
from docx import Document
from docx.shared import Inches, Pt
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt

# Paths
RESULTS_DIR = 'results'
GRAPHS_DIR = os.path.join(RESULTS_DIR, 'graphs')
REPORT_DIR = 'report'

# Ensure report dir exists
os.makedirs(REPORT_DIR, exist_ok=True)

# Load metrics
metrics_path = os.path.join(RESULTS_DIR, 'metrics.csv')
if os.path.exists(metrics_path):
    metrics_df = pd.read_csv(metrics_path)
else:
    metrics_df = pd.DataFrame(columns=["Model", "Accuracy", "Precision", "Recall", "F1-Score", "AUC"])

# Load stats
stats_path = os.path.join(RESULTS_DIR, 'dataset_stats.json')
stats = {}
if os.path.exists(stats_path):
    with open(stats_path, 'r') as f:
        stats = json.load(f)

# --- 1. Generate DOCX (IEEE Format Outline) ---
def generate_docx():
    print("Generating IEEE_Report.docx...")
    doc = Document()
    
    title = doc.add_heading('AI-Powered Fake News Detection Using Text Classification', 0)
    title.alignment = 1
    
    doc.add_heading('Abstract', level=1)
    doc.add_paragraph("This project implements a complete machine learning pipeline to classify news articles as real or fake. "
                      "Various text preprocessing techniques, feature extraction methods (TF-IDF, Bag of Words, Word2Vec), "
                      "and models (KNN, Logistic Regression, Random Forest, Neural Networks) are compared.")
    
    doc.add_heading('1. Introduction', level=1)
    doc.add_paragraph("Fake news poses a significant threat to information integrity. "
                      "This project aims to automate the detection of misinformation using NLP and ML.")
    
    doc.add_heading('2. Dataset Description', level=1)
    doc.add_paragraph("The dataset used is the Kaggle Fake News dataset.")
    if stats:
        doc.add_paragraph("Dataset Statistics:")
        for k, v in stats.items():
            doc.add_paragraph(f"{k}: {v}", style='List Bullet')
    
    doc.add_heading('3. Methodology', level=1)
    doc.add_paragraph("Preprocessing included lowercase conversion, punctuation removal, stopword removal, tokenization, and lemmatization. "
                      "Feature extraction used TF-IDF. Models trained include KNN, Logistic Regression, Random Forest, and MLP Classifier.")
    
    doc.add_heading('4. Results', level=1)
    if not metrics_df.empty:
        doc.add_paragraph("Model Comparison Table:")
        table = doc.add_table(rows=1, cols=len(metrics_df.columns))
        table.style = 'Table Grid'
        hdr_cells = table.rows[0].cells
        for i, col in enumerate(metrics_df.columns):
            hdr_cells[i].text = col
        for _, row in metrics_df.iterrows():
            row_cells = table.add_row().cells
            for i, val in enumerate(row):
                if isinstance(val, float):
                    row_cells[i].text = f"{val:.4f}"
                else:
                    row_cells[i].text = str(val)
                    
    # Graphs
    doc.add_heading('4.1 Model Accuracy & ROC-AUC', level=2)
    acc_graph = os.path.join(GRAPHS_DIR, 'accuracy_comparison.png')
    roc_graph = os.path.join(GRAPHS_DIR, 'roc_curve.png')
    if os.path.exists(acc_graph):
        doc.add_picture(acc_graph, width=Inches(5))
    if os.path.exists(roc_graph):
        doc.add_picture(roc_graph, width=Inches(5))
        
    doc.add_heading('4.2 Confusion Matrices', level=2)
    for model_name in ["KNN", "Logistic Regression", "Random Forest", "Neural Network", "DistilBERT (Fine-Tuned)"]:
        cm_path = os.path.join(GRAPHS_DIR, f'cm_{model_name.replace(" ", "_")}.png')
        if os.path.exists(cm_path):
            doc.add_paragraph(f"{model_name} Confusion Matrix:")
            doc.add_picture(cm_path, width=Inches(3))
        
    doc.add_heading('5. Discussion', level=1)
    doc.add_paragraph("Parametric models like Logistic Regression often perform well and quickly on linearly separable text features (TF-IDF). "
                      "Non-parametric models like KNN suffer from the curse of dimensionality. "
                      "Ensemble methods (Random Forest) and Neural Networks offer high accuracy but take longer to train.")
    
    doc.add_heading('6. Conclusion', level=1)
    doc.add_paragraph("The pipeline successfully identifies fake news. Future scope includes utilizing advanced embeddings like BERT.")
    
    doc.save(os.path.join(REPORT_DIR, 'IEEE_Report.docx'))

# --- 2. Generate PDF (IEEE Format Outline) ---
class IEEEPDF(FPDF):
    def header(self):
        self.set_font('Arial', 'B', 14)
        self.cell(0, 10, 'AI-Powered Fake News Detection Using Text Classification', 0, 1, 'C')
        self.ln(5)

def generate_pdf():
    print("Generating IEEE_Report.pdf...")
    pdf = IEEEPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=11)
    
    def add_section(title, text):
        pdf.set_font("Arial", 'B', 12)
        pdf.cell(0, 10, title, 0, 1)
        pdf.set_font("Arial", '', 11)
        pdf.multi_cell(0, 8, text)
        pdf.ln(5)
        
    add_section("Abstract", "This project implements a complete machine learning pipeline to classify news articles as real or fake. It evaluates KNN, Logistic Regression, Random Forest, and Neural Networks.")
    add_section("1. Introduction", "Fake news poses a significant threat to information integrity.")
    
    stats_text = "The dataset used is the Kaggle Fake News dataset.\n"
    if stats:
        for k, v in stats.items():
            stats_text += f"- {k}: {v}\n"
    add_section("2. Dataset Description", stats_text)
    
    add_section("3. Methodology", "Preprocessing included lowercase conversion, punctuation removal, stopword removal, tokenization, and lemmatization.")
    
    pdf.set_font("Arial", 'B', 12)
    pdf.cell(0, 10, "4. Results", 0, 1)
    pdf.set_font("Arial", '', 9)
    
    if not metrics_df.empty:
        col_width = 190 / len(metrics_df.columns)
        for col in metrics_df.columns:
            pdf.cell(col_width, 10, col, 1)
        pdf.ln()
        for _, row in metrics_df.iterrows():
            for val in row:
                text_val = f"{val:.4f}" if isinstance(val, float) else str(val)
                pdf.cell(col_width, 10, text_val, 1)
            pdf.ln()
            
    acc_graph = os.path.join(GRAPHS_DIR, 'accuracy_comparison.png')
    roc_graph = os.path.join(GRAPHS_DIR, 'roc_curve.png')
    
    pdf.ln(5)
    if os.path.exists(acc_graph):
        pdf.image(acc_graph, w=150)
    pdf.ln(5)
    
    if os.path.exists(roc_graph):
        pdf.image(roc_graph, w=150)
        
    pdf.add_page()
    pdf.set_font("Arial", 'B', 12)
    pdf.cell(0, 10, "Confusion Matrices", 0, 1)
    
    y_pos = pdf.get_y()
    for i, model_name in enumerate(["KNN", "Logistic Regression", "Random Forest", "Neural Network", "DistilBERT (Fine-Tuned)"]):
        cm_path = os.path.join(GRAPHS_DIR, f'cm_{model_name.replace(" ", "_")}.png')
        if os.path.exists(cm_path):
            x = 10 if i % 2 == 0 else 110
            if i % 2 == 0 and i > 0:
                y_pos += 80
            pdf.image(cm_path, x=x, y=y_pos, w=80)
    
    pdf.ln(90)
    pdf.set_font("Arial", 'B', 12)
    pdf.cell(0, 10, "5. Discussion & 6. Conclusion", 0, 1)
    pdf.set_font("Arial", '', 11)
    pdf.multi_cell(0, 8, "Parametric models like Logistic Regression perform robustly on TF-IDF. Future scope includes deep learning embeddings.")
    
    pdf.output(os.path.join(REPORT_DIR, 'IEEE_Report.pdf'))

# --- 3. Generate PPTX ---
def generate_pptx():
    print("Generating Presentation.pptx...")
    prs = Presentation()
    
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "AI-Powered Fake News Detection"
    subtitle.text = "Machine Learning Pipeline\nText Classification"
    
    if stats:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Dataset Statistics"
        tf = slide.placeholders[1].text_frame
        for k, v in stats.items():
            p = tf.add_paragraph()
            p.text = f"{k}: {v}"
    
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Results: Model Accuracy & ROC-AUC"
    acc_graph = os.path.join(GRAPHS_DIR, 'accuracy_comparison.png')
    roc_graph = os.path.join(GRAPHS_DIR, 'roc_curve.png')
    if os.path.exists(acc_graph) and os.path.exists(roc_graph):
        slide.shapes.add_picture(acc_graph, PptxInches(0.5), PptxInches(2), width=PptxInches(4))
        slide.shapes.add_picture(roc_graph, PptxInches(5), PptxInches(2), width=PptxInches(4))
        
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Confusion Matrices"
    
    models_list = ["Logistic Regression", "Random Forest", "DistilBERT (Fine-Tuned)"]
    for i, model_name in enumerate(models_list):
        cm_path = os.path.join(GRAPHS_DIR, f'cm_{model_name.replace(" ", "_")}.png')
        if os.path.exists(cm_path):
            x = PptxInches(0.5 + i * 3.0)
            slide.shapes.add_picture(cm_path, x, PptxInches(2), width=PptxInches(2.8))
        
    prs.save(os.path.join(REPORT_DIR, 'Presentation.pptx'))

if __name__ == "__main__":
    generate_docx()
    generate_pdf()
    generate_pptx()
    print("All reports generated successfully in the 'report' directory.")
